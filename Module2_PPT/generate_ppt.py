from pptx import Presentation
from pptx.util import Inches, Pt
import glob
import os

prs = Presentation()

# Slide 1: Title Slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Module 2: The Data Lab"
subtitle.text = "Transforming Messy Data into Actionable Insights"

# Define a function to add standard slide
def add_slide(prs, headline, bullet_points, image_path=None):
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = headline
    tf = body_shape.text_frame
    
    for i, bp in enumerate(bullet_points):
        if i == 0:
            tf.text = bp
        else:
            p = tf.add_paragraph()
            p.text = bp
            
    if image_path and os.path.exists(image_path):
        # Resize/reposition the text frame if an image is present
        body_shape.width = Inches(4.5)
        slide.shapes.add_picture(image_path, Inches(5), Inches(2), width=Inches(4.5))

# Slide 2
image2 = glob.glob("messy_data_concrete*.png")
add_slide(prs, "Data is the New Concrete", 
          ["Every construction, agricultural, and industrial project runs on data.", 
           "The Challenge: Finding anomalies in thousands of rows manually takes hours.", 
           "The Solution: AI Data Analysts clean and organize data in seconds.",
           "Activity: How long would it take you to manually find the 5 budget overruns here?"], 
          image2[0] if image2 else None)

# Slide 3
add_slide(prs, "Welcome to the Data Lab", 
          ["Step 1: Upload your dataset (e.g., procurement_demo.csv).", 
           "Step 2: Write your first AI cleaning prompt.", 
           "Step 3: Watch the AI fix inconsistent vendor names and messy dates instantly!"], 
          None)

# Slide 4
image4 = glob.glob("invoice_reconciliation*.png")
add_slide(prs, "The Great Reconciliation Challenge", 
          ["Scenario: Your Invoice doesn't match the Goods Received Note.", 
           "The Mission: Use the AI Reconciliation Engine to find the missing ₹5 Lakhs.", 
           "The Rules: The first 3 people to find the discrepancy win!"], 
          image4[0] if image4 else None)

# Slide 5
image5 = glob.glob("mis_generator_dashboard*.png")
add_slide(prs, "From Data to Visual Storytelling", 
          ["Clean data is only half the battle. Presenting it is the other half.", 
           "Generate automated visual reports for management.", 
           "Instantly create comparison charts (e.g., Expenditure vs. Budgeted BOQ)."], 
          image5[0] if image5 else None)

# Slide 6
image6 = glob.glob("dec_industries_breakout*.png")
add_slide(prs, "Real-World Application by Department", 
          ["DEC Infra: Analyze Safety Audits & Construction Timelines.", 
           "DEC Agro: Analyze Cold Chain Temperatures & Crop Yield Data.", 
           "DEC Industries: Analyze Equipment Downtime & Manufacturing Efficiency."], 
          image6[0] if image6 else None)

# Slide 7
add_slide(prs, "Key Takeaways & What's Next", 
          ["What was your biggest surprise today?", 
           "How will this change your workflow tomorrow?", 
           "Coming Up Next: Module 3 - Safe AI Usage (Security & Ethics)."], 
          None)

prs.save('Module2_Engagement_Plan.pptx')
print("PPTX generated successfully.")
